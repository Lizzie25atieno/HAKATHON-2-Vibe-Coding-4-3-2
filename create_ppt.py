from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# Define colors and fonts
HEADER_COLOR = RGBColor(44, 62, 80)       # Dark blue
CONTENT_COLOR = RGBColor(33, 37, 41)      # Dark gray
BG_COLOR = RGBColor(244, 246, 249)        # Light gray
ACCENT_COLOR = RGBColor(52, 152, 219)     # Blue accent for icons/bullets
FONT_NAME = 'Arial'

# Define slides content
slides_content = [
    ("Problem Clarity", 
     "• Finding study partners is challenging.\n"
     "• Limited access to lecturers and exercises.\n"
     "• No centralized platform for internships and developer tools."),
    
    ("Solution Quality", 
     "• Connect students by course, skills, and university.\n"
     "• Real-time chat, buddy requests, and skill-based filtering.\n"
     "• Access lecturers, exercises, internships, and developer tools in one platform."),
    
    ("Market Insight", 
     "• Target: University students aged 18–25.\n"
     "• One-stop platform improves collaboration, mentorship, skill development.\n"
     "• Enhances student engagement and networking opportunities.")
]

# Add slides
for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content layout
    
    # Slide background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Slide title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.name = FONT_NAME
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = HEADER_COLOR
    
    # Add white rounded rectangle behind content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    rect.line.color.rgb = RGBColor(200, 200, 200)
    
    # Content text
    content_box = slide.placeholders[1]
    content_box.text = content
    for para in content_box.text_frame.paragraphs:
        para.font.name = FONT_NAME
        para.font.size = Pt(24)
        para.font.color.rgb = CONTENT_COLOR
        para.space_after = Pt(10)
        para.alignment = PP_ALIGN.LEFT
        
        # Make bullets accent color
        if para.text.startswith("•"):
            para.level = 0
            para.font.color.rgb = ACCENT_COLOR
    
    # Optional: Add small icon circles to each bullet
    lines = content.split("\n")
    for i, line in enumerate(lines):
        if line.startswith("•"):
            top_icon = top + Inches(0.25) + i * Pt(30).pt / 72  # approximate vertical spacing
            slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top + 0.4 + i*0.7, Inches(0.2), Inches(0.2)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = ACCENT_COLOR

# Save presentation
file_path = "StudyBuddy_Pitch.pptx"
prs.save(file_path)
print(f'Presentation saved to {file_path}')
